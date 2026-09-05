#!/usr/bin/env python3
"""粗检:文本溢出(按 CJK 估宽)、出界、文本框重叠"""
import math, sys
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
prs = Presentation("/Users/xierenfeng/Projects/ai/weknora/agent-intro-ppt/WeKnora智能体介绍.pptx")
W, H = prs.slide_width / EMU, prs.slide_height / EMU
issues = 0

def char_w(ch, pt):
    return pt * (0.56 if ord(ch) < 256 else 1.0)

for si, slide in enumerate(prs.slides, 1):
    boxes = []
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        x, y = sh.left / EMU, sh.top / EMU
        w, h = sh.width / EMU, sh.height / EMU
        if x < -0.05 or y < -0.05 or x + w > W + 0.05 or y + h > H + 0.05:
            print(f"S{si} OUT-OF-SLIDE: '{sh.text_frame.text[:14]}' at ({x:.2f},{y:.2f},{w:.2f},{h:.2f})")
            issues += 1
        # 估算所需高度
        need = 0.0
        for p in sh.text_frame.paragraphs:
            pt = 12
            for r in p.runs:
                if r.font.size: pt = r.font.size.pt; break
            line_w = sum(char_w(c, pt) for c in p.text) / 72.0
            lines = max(1, math.ceil(line_w / max(w, 0.1))) if p.text else 1
            need += lines * pt * 1.35 / 72.0
        if need > h + 0.12:
            print(f"S{si} OVERFLOW?: '{sh.text_frame.text[:18]}' need~{need:.2f} box_h={h:.2f} w={w:.2f}")
            issues += 1
        boxes.append((x, y, w, h, sh.text_frame.text[:12]))
    # 重叠检测(仅报告文本框间显著重叠)
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            ox = min(a[0]+a[2], b[0]+b[2]) - max(a[0], b[0])
            oy = min(a[1]+a[3], b[1]+b[3]) - max(a[1], b[1])
            if ox > 0.15 and oy > 0.15:
                print(f"S{si} OVERLAP: '{a[4]}' vs '{b[4]}' ({ox:.2f}x{oy:.2f})")
                issues += 1
print(f"\ntotal potential issues: {issues}")
